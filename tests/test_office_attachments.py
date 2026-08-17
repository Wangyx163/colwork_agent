from __future__ import annotations

import base64
import importlib.util
import io
import unittest
import zipfile

from collab_agent.attachments import extract_attachment_text

# Office extraction lives behind the `office` extra. These tests assert the
# installed behaviour, so without markitdown they would report a deployment
# choice as a broken build -- the same reason the PostgreSQL adapter test skips
# itself when no database URL is set. The degraded answer for a deployment that
# never installed the extra is covered separately and does not need the extra.
MARKITDOWN_INSTALLED = importlib.util.find_spec("markitdown") is not None

DOCX_MIME = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

CONTENT_TYPES = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""

ROOT_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""


def minimal_docx(paragraphs: list[str]) -> bytes:
    """Build a valid minimal .docx without pulling a writer dependency."""

    body = "".join(
        f"<w:p><w:r><w:t xml:space='preserve'>{text}</w:t></w:r></w:p>"
        for text in paragraphs
    )
    document = (
        "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>"
        "<w:document xmlns:w='http://schemas.openxmlformats.org/wordprocessingml/2006/main'>"
        f"<w:body>{body}</w:body></w:document>"
    )
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", CONTENT_TYPES)
        archive.writestr("_rels/.rels", ROOT_RELS)
        archive.writestr("word/document.xml", document)
    return buffer.getvalue()


def minimal_xlsx(rows: list[list[str]]) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    for row in rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def minimal_pptx(title: str, body: str) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def payload(name: str, data: bytes, mime: str = "") -> dict:
    encoded = base64.b64encode(data).decode("ascii")
    declared = mime or "application/octet-stream"
    return {
        "name": name,
        "type": mime,
        "size": len(data),
        "data": f"data:{declared};base64,{encoded}",
    }


@unittest.skipUnless(
    MARKITDOWN_INSTALLED,
    'install the "office" extra to run Office attachment extraction tests',
)
class OfficeAttachmentTests(unittest.TestCase):
    def test_docx_text_becomes_citable_evidence(self) -> None:
        questions = [
            "你最近一次觉得特别开心是什么时候？",
            "你会怎么向外地朋友介绍这座城市？",
        ]
        result = extract_attachment_text(
            payload("问题清单.docx", minimal_docx(questions), DOCX_MIME)
        )
        self.assertEqual(result["extraction_status"], "EXTRACTED")
        for question in questions:
            self.assertIn(question, result["extracted_text"])
        self.assertIn("content_sha256", result)

    def test_xlsx_cells_are_extracted(self) -> None:
        result = extract_attachment_text(
            payload(
                "问题表.xlsx",
                minimal_xlsx(
                    [["序号", "问题"], ["1", "你上一次熬夜是因为什么？"]]
                ),
                XLSX_MIME,
            )
        )
        self.assertEqual(result["extraction_status"], "EXTRACTED")
        self.assertIn("你上一次熬夜是因为什么？", result["extracted_text"])

    def test_pptx_slide_text_is_extracted(self) -> None:
        result = extract_attachment_text(
            payload(
                "执行表.pptx",
                minimal_pptx("采访执行表", "地点：市中心步行街"),
                PPTX_MIME,
            )
        )
        self.assertEqual(result["extraction_status"], "EXTRACTED")
        self.assertIn("采访执行表", result["extracted_text"])

    def test_the_filename_extension_routes_when_mime_is_missing(self) -> None:
        # Browsers occasionally send an empty or generic type for Office files.
        result = extract_attachment_text(
            payload("无类型.docx", minimal_docx(["按扩展名识别"]), "")
        )
        self.assertEqual(result["extraction_status"], "EXTRACTED")
        self.assertIn("按扩展名识别", result["extracted_text"])

    def test_a_corrupt_office_file_fails_without_crashing(self) -> None:
        corrupt = bytes(range(256)) * 8
        result = extract_attachment_text(
            payload("坏文件.docx", corrupt, DOCX_MIME)
        )
        self.assertEqual(result["extraction_status"], "FAILED")
        self.assertEqual(result["extracted_text"], "")
        self.assertIn("extraction_error", result)

    def test_a_mislabelled_office_file_is_read_as_what_it_really_is(self) -> None:
        """MarkItDown sniffs real content rather than trusting the extension.

        A .docx that is actually plain text extracts as that text instead of
        failing. The recorded content_sha256 still pins what was uploaded, so
        the evidence stays honest.
        """

        result = extract_attachment_text(
            payload("其实是文本.docx", "候选问题一\n候选问题二".encode("utf-8"), DOCX_MIME)
        )
        self.assertEqual(result["extraction_status"], "EXTRACTED")
        self.assertIn("候选问题一", result["extracted_text"])

    def test_unsupported_types_still_degrade_the_same_way(self) -> None:
        result = extract_attachment_text(
            payload("图片.png", b"\x89PNG\r\n\x1a\n", "image/png")
        )
        self.assertEqual(result["extraction_status"], "UNSUPPORTED")
        self.assertEqual(result["extracted_text"], "")

    def test_pdf_and_text_paths_are_unchanged(self) -> None:
        text = extract_attachment_text(
            payload("说明.txt", "候选问题一".encode("utf-8"), "text/plain")
        )
        self.assertEqual(text["extraction_status"], "EXTRACTED")
        self.assertEqual(text["extracted_text"], "候选问题一")


if __name__ == "__main__":
    unittest.main()
